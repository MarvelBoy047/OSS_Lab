from pathlib import Path
import fitz  # PyMuPDF for PDFs
import docx
from pptx import Presentation

def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix == ".docx":
        return extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    elif suffix in [".txt"]:
        return extract_text_from_txt(file_path)
    else:
        return ""

def extract_text_from_pdf(file_path: Path) -> str:
    text = ""
    try:
        doc = fitz.open(str(file_path))
        for page in doc:
            text += page.get_text()
        return text
    except Exception:
        return ""

def extract_text_from_docx(file_path: Path) -> str:
    try:
        doc = docx.Document(str(file_path))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception:
        return ""

def extract_text_from_pptx(file_path: Path) -> str:
    text = []
    try:
        pres = Presentation(str(file_path))
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        return ""

def extract_text_from_txt(file_path: Path) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception:
        return ""
